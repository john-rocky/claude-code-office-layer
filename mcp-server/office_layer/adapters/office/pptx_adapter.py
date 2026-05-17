"""python-pptx adapter — PowerPoint .pptx."""

from __future__ import annotations

from pathlib import Path

from ...models import ChunkKind, DocumentChunk, DocumentKind, ExtractedTable, ExtractionMethod
from ..base import ExtractionResult


class PythonPptxAdapter:
    name = "python-pptx"
    supported_kinds = (DocumentKind.PPTX,)

    def __init__(self) -> None:
        try:
            import pptx  # noqa: F401
            self._available = True
        except ImportError:
            self._available = False

    def is_available(self) -> bool:
        return self._available

    def extract(
        self,
        path: Path,
        kind: DocumentKind,
        *,
        document_id: str,
        workspace_id: str,
    ) -> ExtractionResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            return ExtractionResult(error=f"python-pptx missing: {exc}")

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            return ExtractionResult(error=f"python-pptx failed to open: {exc}")

        chunks: list[DocumentChunk] = []
        tables: list[ExtractedTable] = []
        ordinal = 0
        slide_count = 0

        for i, slide in enumerate(prs.slides, start=1):
            slide_count += 1
            title = None
            body_parts: list[str] = []
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text.strip() or None
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text and text != title:
                            body_parts.append(text)
                if shape.has_table:
                    rows = []
                    for r in shape.table.rows:
                        rows.append([cell.text.strip() for cell in r.cells])
                    if rows:
                        headers = rows[0]
                        body = rows[1:] if len(rows) > 1 else []
                        tables.append(
                            ExtractedTable(
                                document_id=document_id,
                                slide_number=i if False else None,  # spec doesn't track slide on table
                                headers=headers,
                                rows=body,
                                confidence=0.85,
                            )
                        )
            # Notes
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except AttributeError:
                notes_text = ""
            if notes_text:
                body_parts.append(f"[notes]\n{notes_text}")

            slide_text = "\n".join([t for t in [title, *body_parts] if t]).strip()
            if not slide_text:
                continue
            chunks.append(
                DocumentChunk(
                    id=DocumentChunk.make_id(document_id, ordinal),
                    document_id=document_id,
                    workspace_id=workspace_id,
                    kind=ChunkKind.PPTX_SLIDE,
                    ordinal=ordinal,
                    slide_number=i,
                    section_heading=title,
                    text=slide_text,
                    char_count=len(slide_text),
                    confidence=1.0,
                    extraction_method=ExtractionMethod.PPTX_PARSE,
                )
            )
            ordinal += 1

        return ExtractionResult(
            chunks=chunks,
            tables=tables,
            slide_count=slide_count or None,
            extraction_method=ExtractionMethod.PPTX_PARSE,
        )
